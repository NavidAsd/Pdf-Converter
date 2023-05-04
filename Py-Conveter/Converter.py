from pdf2docx import Converter
from dataclasses import replace
from fileinput import filename
import os
import win32com.client
import pandas 
import pdfkit
from pdf2image import convert_from_path
from pptx import Presentation
from io import BytesIO
import pythoncom
import fitz
from AppliedMethodes import Methodes
import time

class Program():
    def __init__():
        pass

    def PptToPdf(InputFilePath,InputFileName,OutputPath,UserIp):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None):
            try:
                outFileName= Methodes.ReturnFileName('.pdf','Pdf')
                OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\{UserIp}'
                InputFilePath = Methodes.FixRequestFormat(InputFilePath)
                Methodes.CreateDirectory(OutputPath)

                in_file=f'{InputFilePath}\{InputFileName}'
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application",pythoncom.CoInitialize())
                deck = powerpoint.Presentations.Open(in_file)
                
                time.sleep(1)
                deck.SaveAs(f'{OutputPath}\\{outFileName}', 32) # formatType = 32 for ppt to pdf
                
                return {'Success':True,'Message':'ConvertSuccessFully','OutFile':outFileName,'OutPath':Methodes.ReverseRequestFormat(OutputPath)} 
            except:
                return {'Success':False,'Message':'Error In Converting Proccess'}
            finally:
                try:
                    os.system("TASKKILL /F /IM powerpnt.exe")
                    deck.Close()
                except:
                    pass
                try:
                    powerpoint.Quit()
                except:
                    pass
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}
    
    def ExcelToPdf(InputFilePath,InputFileName,OutputPath,UserIp):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None):
            try:
                outFileName= Methodes.ReturnFileName('.pdf','Pdf')
                OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\\{UserIp}'
                InputFilePath = Methodes.FixRequestFormat(InputFilePath)
                Methodes.CreateDirectory(OutputPath)

                df = pandas.read_excel(f'{InputFilePath}\\{InputFileName}')
                df.to_html(f"file.html")
                config = pdfkit.configuration(wkhtmltopdf='E:/C#/Web/Test/ConvertFromPdf/wkhtmltopdf/bin/wkhtmltopdf.exe')
                pdfkit.from_file("file.html", f"{OutputPath}\\{outFileName}",configuration=config)
                os.remove(os.path.join("file.html"))
                return {'Success':True,'Message':'ConvertSuccessFuly','OutFile':outFileName,'OutPath':Methodes.ReverseRequestFormat(OutputPath)} 
            except:
                return {'Success':False,'Message':'Error In Converting Proccess'}
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}

    def PdfToPpt(InputFilePath,InputFileName,OutputPath,UserIp):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None):
            
            outFileName= Methodes.ReturnFileName('.pptx','PowerPoint')
            OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\\{UserIp}'
            InputFilePath = Methodes.FixRequestFormat(InputFilePath)
            Methodes.CreateDirectory(OutputPath)

            # Prep presentation
            try:
                prs = Presentation()
                blank_slide_layout = prs.slide_layouts[6]

                # Convert PDF to list of images
                #print("Starting conversion...")
                slideimgs = convert_from_path(f'{InputFilePath}/{InputFileName}', 300, fmt='ppm', thread_count=2,poppler_path='./poppler-0.68.0/bin')
                #print("...complete.")

                # Loop over slides
                for i, slideimg in enumerate(slideimgs):
                    #if i % 10 == 0:
                        #print("Saving slide: " + str(i))

                    imagefile = BytesIO()
                    slideimg.save(imagefile, format='tiff')
                    #imagefile.getvalue()
                    imagefile.seek(0)
                    width, height = slideimg.size

                    # Set slide dimensions
                    prs.slide_height = height * 9525
                    prs.slide_width = width * 9525

                    # Add slide
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

                # Save Powerpoint
                #print("Saving file: " + base_name + ".pptx")
                prs.save(f'{OutputPath}/{outFileName}')
                #print("Conversion complete. :)")
                return {'Success':True,'Message':'ConvertSuccessFully','OutFile':outFileName,'OutPath':Methodes.ReverseRequestFormat(OutputPath)} 
            except:
                return {'Success':False,'Message':'Error In Converting Proccess'}
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}

    def WordToPdf(InputFilePath,InputFileName,OutputPath,UserIp):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None):
            try:
                outFileName= Methodes.ReturnFileName('.pdf','Pdf')
                OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\\{UserIp}'
                InputFilePath = Methodes.FixRequestFormat(InputFilePath)

                in_file = f'{InputFilePath}\\{InputFileName}'
                out_file = f'{OutputPath}\\{outFileName}'

                Methodes.CreateDirectory(OutputPath)
                word = win32com.client.Dispatch('Word.Application',pythoncom.CoInitialize())
                doc = word.Documents.Open(in_file)
                doc.SaveAs(out_file, FileFormat = 17) # formatType = 17 for word to pdf
                try:
                    doc.Close()
                except:
                    pass
                #try:
                    #word.Quit()
                #except:
                    #pass
                return {'Success':True,'Message':'ConvertSuccessFuly','OutFile':outFileName,'OutPath':Methodes.ReverseRequestFormat(OutputPath)}
            except:
                return {'Success':False,'Message':'Error In Converting Proccess'}
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}

    def PdfToWord(InputFilePath,InputFileName,OutputPath,UserIp,FormatType):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None and int(FormatType) >= 0):
            try:

                outFileName= Methodes.ReturnFileName(Methodes.SetDocFormat(int(FormatType)),'Word')
                OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\\{UserIp}'
                InputFilePath = Methodes.FixRequestFormat(InputFilePath)   

                in_file = f'{InputFilePath}\\{InputFileName}'
                out_file = f'{OutputPath}\\{outFileName}'

                Methodes.CreateDirectory(OutputPath)

                cv =Converter(in_file)
                cv.convert(out_file,start=0,end=None)
                cv.close()

                if(os.path.isfile(out_file) == False):
                    return {'Success':False,'Message':'Error In Converting Proccess'}

                return {'Success':True,'Message':'ConvertSuccessFuly','OutFile':outFileName,'OutPath':Methodes.ReverseRequestFormat(OutputPath)}
            except:
                return {'Success':False,'Message':'Error In Converting Proccess'}
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}

    def ExtractPdfImages(InputFilePath,InputFileName,OutputPath,UserIp):
        if(InputFilePath,InputFileName,OutputPath,UserIp != None):
            outfilename = Methodes.ReturnFileNameWithoutModifie('.zip','ExtractImages_Compressed')
            OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\\{UserIp}'
            InputFilePath = Methodes.FixRequestFormat(InputFilePath)
            Methodes.CreateDirectory(OutputPath)

            try:
                doc = fitz.open(f"{InputFilePath}\\{InputFileName}")
                count=0
                files=[]
                for i in range(len(doc)):
                    for img in doc.getPageImageList(i):
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:       # this is GRAY or RGB
                            count+=1
                            file =f"{OutputPath}\\{Methodes.ReturnFileNameWithoutDate('.png',f'ExtractImages-{xref}',count)}"
                            files.append(file)
                            pix.writePNG(file)
                        else:               # CMYK: convert to RGB first
                            count+=1
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            pix1.writePNG("p%s-%s.png" % (i, xref))
                            pix1 = None
                        pix = None
                result = Methodes.CompressFilesToZip(files,OutputPath,outfilename)
                for i in files:
                    try:
                        os.remove(i)
                    except:
                        pass
                if(result):
                    return {'Success':True,'Message':'ConvertSuccessFuly','OutFile':outfilename,'OutPath':Methodes.ReverseRequestFormat(OutputPath)}
                else:
                    return {'Success':False,'Message':'Error In Extracting Proccess'}
            except:
                return {'Success':False,'Message':'Error In Extracting Proccess'}
        else:
            return {'Success':False,'Message':'Error Input Cannot Be null'}


#Program.WordToPdf('E:\\C\\Input','test2.pdf','E:\\C\\Output','192.168.1.103')
#Program.ExtractPdfImages('E:\\C\\Input','Test-English.pdf','E:\\C\\Output','192.168.1.103')

# UserIp='127.0.0.1'
# OutputPath='E:-C-Output'
# InputFilePath='E:-C-Input'
# InputFileName='LinuxArticle.docx'

# outFileName= Methodes.ReturnFileName('.pdf','Pdf')
# OutputPath = f'{Methodes.FixRequestFormat(OutputPath)}\{UserIp}'
# InputFilePath = Methodes.FixRequestFormat(InputFilePath)
# # try:
# #     if(os.path.isdir(OutputPath) == False):
# #         os.mkdir(OutputPath)
# # except:
# #     pass
# in_file=f'{InputFilePath}\{InputFileName}'

# powerpoint = win32com.client.Dispatch("PowerPoint.Application",pythoncom.CoInitialize())
# deck = powerpoint.Presentations.Open(in_file)
# import time
# time.sleep(2)
# deck.SaveAs(f'{OutputPath}\\{outFileName}', 32) # formatType = 32 for ppt to pdfc
# os.system("TASKKILL /F /IM powerpnt.exe")


